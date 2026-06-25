"""从 pptx 课件 / pdf 资料抽取内嵌图片，供视觉识别（VLM）转写。

课件里的图标、流程图、产品对比表、话术截图等往往承载文字之外的关键信息，
需要交给多模态模型转写。本模块只负责「取图 + 预过滤」，不做任何模型调用：

- 按页顺序遍历图片（pptx 含 group 内嵌套图片）；
- 用 sha256 去重，避免同一张图（如每页页脚 logo）被反复识别；
- 用像素尺寸预过滤掉过小的装饰元素（分隔线、小图标），先省一道模型成本；
- 只保留多模态模型可识别的常见位图格式，矢量图 / 异常格式跳过。
"""

from __future__ import annotations

import base64
import hashlib
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader

# 过小的图（最长边 < 该阈值像素）通常是分隔线、项目符号、小图标等装饰元素，
# 识别价值低且浪费 token，直接预过滤。
DEFAULT_MIN_SIDE_PX = 128

# 多模态模型可识别的常见位图格式。矢量图（wmf/emf）与异常格式不传给模型。
_SUPPORTED_MEDIA_TYPES = frozenset(
    {"image/png", "image/jpeg", "image/gif", "image/webp"}
)
# 个别课件里 jpg 的 content_type 写成 image/jpg，归一到标准 image/jpeg。
_MEDIA_TYPE_ALIASES = {"image/jpg": "image/jpeg"}
# PIL 解码出的格式名 → MIME 类型（用于 pdf 抽图）。
_PIL_FORMAT_TO_MEDIA = {
    "JPEG": "image/jpeg",
    "PNG": "image/png",
    "GIF": "image/gif",
    "WEBP": "image/webp",
}


@dataclass(frozen=True)
class SourceImage:
    """素材文件（pptx/pdf）中的一张内嵌图片（已去重、已通过尺寸预过滤）。"""

    page_index: int  # 1-based，对应「第 N 页」
    media_type: str  # 归一后的 MIME 类型，如 image/png
    blob: bytes
    width: int  # 像素；0 表示尺寸未知
    height: int
    sha256: str

    @property
    def max_side(self) -> int:
        """最长边像素；尺寸未知时为 0。"""
        return max(self.width, self.height)

    @property
    def base64_data(self) -> str:
        """base64 编码的图片数据，供构造多模态消息块。"""
        return base64.b64encode(self.blob).decode("ascii")


# ---- 通用去重 + 尺寸预过滤 ----


def _dedup_and_filter(
    raw: list[tuple[int, str, bytes, int, int]], min_side_px: int
) -> tuple[SourceImage, ...]:
    """对 (page_index, media_type, blob, w, h) 原始列表做去重与尺寸预过滤。"""
    seen: set[str] = set()
    images: list[SourceImage] = []
    for page_index, media_type, blob, width, height in raw:
        if media_type not in _SUPPORTED_MEDIA_TYPES:
            continue
        digest = hashlib.sha256(blob).hexdigest()
        if digest in seen:  # 同图（如统一页脚 logo）只识别一次
            continue
        if 0 < max(width, height) < min_side_px:  # 尺寸已知且过小 → 装饰元素
            continue
        seen.add(digest)
        images.append(
            SourceImage(
                page_index=page_index,
                media_type=media_type,
                blob=blob,
                width=width,
                height=height,
                sha256=digest,
            )
        )
    return tuple(images)


# ---- pptx ----


def _iter_picture_shapes(shapes):
    """按顺序产出所有图片形状，递归展开 group。"""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_picture_shapes(shape.shapes)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape


def _pptx_image_size(image) -> tuple[int, int]:
    """读取 pptx 图片像素尺寸；取不到时返回 (0, 0)。"""
    try:
        width, height = image.size
        return int(width), int(height)
    except Exception:  # noqa: BLE001 - 矢量图/异常格式拿不到尺寸，按未知处理
        return 0, 0


def extract_pptx_images(
    path: Path, min_side_px: int = DEFAULT_MIN_SIDE_PX
) -> tuple[SourceImage, ...]:
    """抽取 pptx 内嵌图片，按页序返回去重且通过尺寸预过滤的图片。

    解析失败时返回空元组（与解析层「失败不抛、降级返回」的风格一致）。
    """
    try:
        prs = Presentation(str(path))
    except Exception:  # noqa: BLE001 - 损坏文件不应中断整条流水线
        return ()

    raw: list[tuple[int, str, bytes, int, int]] = []
    for page_index, slide in enumerate(prs.slides, start=1):
        for shape in _iter_picture_shapes(slide.shapes):
            try:
                image = shape.image
            except Exception:  # noqa: BLE001 - 个别图片对象不可读，跳过
                continue
            media_type = _MEDIA_TYPE_ALIASES.get(
                image.content_type, image.content_type
            )
            width, height = _pptx_image_size(image)
            raw.append((page_index, media_type, image.blob, width, height))
    return _dedup_and_filter(raw, min_side_px)


# ---- pdf ----


def extract_pdf_images(
    path: Path, min_side_px: int = DEFAULT_MIN_SIDE_PX
) -> tuple[SourceImage, ...]:
    """抽取 pdf 各页内嵌图片，按页序返回去重且通过尺寸预过滤的图片。

    为未来「以图为主、缺 word」的上传数据预留；解析失败降级返回空元组。
    """
    try:
        reader = PdfReader(str(path))
    except Exception:  # noqa: BLE001 - 损坏文件不应中断整条流水线
        return ()

    raw: list[tuple[int, str, bytes, int, int]] = []
    for page_index, page in enumerate(reader.pages, start=1):
        try:
            page_images = list(page.images)
        except Exception:  # noqa: BLE001 - 个别页图片解析失败，跳过该页
            continue
        for image_file in page_images:
            try:
                pil_image = image_file.image
                media_type = _PIL_FORMAT_TO_MEDIA.get(pil_image.format or "", "")
                width, height = pil_image.size
            except Exception:  # noqa: BLE001 - 单图解码失败，跳过
                continue
            if not media_type:
                continue
            raw.append(
                (page_index, media_type, image_file.data, int(width), int(height))
            )
    return _dedup_and_filter(raw, min_side_px)


# ---- 按扩展名分派 ----


def extract_images(
    path: Path, min_side_px: int = DEFAULT_MIN_SIDE_PX
) -> tuple[SourceImage, ...]:
    """根据文件扩展名抽取内嵌图片（支持 pptx / pdf，其余返回空）。"""
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return extract_pptx_images(path, min_side_px)
    if suffix == ".pdf":
        return extract_pdf_images(path, min_side_px)
    return ()
