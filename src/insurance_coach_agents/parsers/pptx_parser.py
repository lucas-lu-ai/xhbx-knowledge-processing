"""解析 pptx 课件，按幻灯片提取标题与要点文本。

课件提供框架与话术骨架。每页输出 ``## 第 N 页`` 小标题 + 该页所有文本框内容，
保留页边界，便于下游理解章节结构。
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from ..models import FileType, ParsedFile


def _slide_text(slide) -> list[str]:
    """提取单页所有形状的文本（含表格单元格），按出现顺序。"""
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)
        elif shape.has_table:
            for row in shape.table.rows:
                cells = [c.text.strip() for c in row.cells]
                line = " | ".join(c for c in cells if c)
                if line:
                    parts.append(line)
    return parts


def parse_pptx(path: Path) -> ParsedFile:
    """解析 pptx 文件为 ``ParsedFile``。"""
    filename = path.name
    try:
        prs = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        return ParsedFile(
            file_type=FileType.PPTX,
            filename=filename,
            text="",
            warnings=(f"pptx 解析失败: {exc}",),
        )

    blocks: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        parts = _slide_text(slide)
        if parts:
            blocks.append(f"## 第 {index} 页\n" + "\n".join(parts))

    warnings = () if blocks else ("pptx 无文本内容",)
    return ParsedFile(
        file_type=FileType.PPTX,
        filename=filename,
        text="\n\n".join(blocks),
        warnings=warnings,
    )
