"""视觉增强相关单元测试：图片抽取 / 识别 / 章节组装（fake 模型，不调真实 API）。"""

from __future__ import annotations

import asyncio
import base64
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from insurance_coach_agents.agents.enrich import enrich_section_with_vision
from insurance_coach_agents.agents.vision import ImageDescriber
from insurance_coach_agents.parsers.image_extract import (
    SourceImage,
    extract_images,
    extract_pdf_images,
    extract_pptx_images,
)

# 合成测试图：两张 160x160（足够大，不被默认尺寸预过滤）、一张 16x16（应被过滤）。
_BIG_RED = "iVBORw0KGgoAAAANSUhEUgAAAKAAAACgCAIAAAAErfB6AAABiUlEQVR4nO3RAQkAIADAMDWD/ZMZxhQinC3B4fPsPehavwN4y+A4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g0fbBX7SAkSR6z3PAAAAAElFTkSuQmCC"  # noqa: E501
_BIG_BLUE = "iVBORw0KGgoAAAANSUhEUgAAAKAAAACgCAIAAAAErfB6AAABiUlEQVR4nO3RAQkAIADAMDWD/ZMZxhQinC3B4XPvM+havwN4y+A4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g+MMjjM4zuA4g0fbBX1CAkTCK4xqAAAAAElFTkSuQmCC"  # noqa: E501
_SMALL = "iVBORw0KGgoAAAANSUhEUgAAABAAAAAQCAIAAACQkWg2AAAAIklEQVR4nGNkOMFAEmAiTTnDqAaiABNxyhBgVAMxgORQAgCXHgDoLRdxkQAAAABJRU5ErkJggg=="  # noqa: E501
# 含一张 200x200 红图的单页 PDF（Pillow PDF 驱动生成），用于验证 pdf 抽图。
_PDF_WITH_IMAGE = "JVBERi0xLjQKJSBjcmVhdGVkIGJ5IFBpbGxvdyBQREYgZHJpdmVyCjQgMCBvYmo8PAovVHlwZSAvQ2F0YWxvZwovUGFnZXMgNSAwIFIKPj5lbmRvYmoKNSAwIG9iajw8Ci9UeXBlIC9QYWdlcwovQ291bnQgMQovS2lkcyBbIDIgMCBSIF0KPj5lbmRvYmoKMSAwIG9iajw8Ci9UeXBlIC9YT2JqZWN0Ci9TdWJ0eXBlIC9JbWFnZQovV2lkdGggMjAwCi9IZWlnaHQgMjAwCi9GaWx0ZXIgL0RDVERlY29kZQovQml0c1BlckNvbXBvbmVudCA4Ci9Db2xvclNwYWNlIC9EZXZpY2VSR0IKL0xlbmd0aCAxMzA1Cj4+c3RyZWFtCv/Y/+AAEEpGSUYAAQEAAAEAAQAA/9sAQwAIBgYHBgUIBwcHCQkICgwUDQwLCwwZEhMPFB0aHx4dGhwcICQuJyAiLCMcHCg3KSwwMTQ0NB8nOT04MjwuMzQy/9sAQwEJCQkMCwwYDQ0YMiEcITIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIy/8AAEQgAyADIAwEiAAIRAQMRAf/EAB8AAAEFAQEBAQEBAAAAAAAAAAABAgMEBQYHCAkKC//EALUQAAIBAwMCBAMFBQQEAAABfQECAwAEEQUSITFBBhNRYQcicRQygZGhCCNCscEVUtHwJDNicoIJChYXGBkaJSYnKCkqNDU2Nzg5OkNERUZHSElKU1RVVldYWVpjZGVmZ2hpanN0dXZ3eHl6g4SFhoeIiYqSk5SVlpeYmZqio6Slpqeoqaqys7S1tre4ubrCw8TFxsfIycrS09TV1tfY2drh4uPk5ebn6Onq8fLz9PX29/j5+v/EAB8BAAMBAQEBAQEBAQEAAAAAAAABAgMEBQYHCAkKC//EALURAAIBAgQEAwQHBQQEAAECdwABAgMRBAUhMQYSQVEHYXETIjKBCBRCkaGxwQkjM1LwFWJy0QoWJDThJfEXGBkaJicoKSo1Njc4OTpDREVGR0hJSlNUVVZXWFlaY2RlZmdoaWpzdHV2d3h5eoKDhIWGh4iJipKTlJWWl5iZmqKjpKWmp6ipqrKztLW2t7i5usLDxMXGx8jJytLT1NXW19jZ2uLj5OXm5+jp6vLz9PX29/j5+v/aAAwDAQACEQMRAD8A5CiiivFP0wKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooAKKKKACiiigAooooA//2QplbmRzdHJlYW0KZW5kb2JqCjIgMCBvYmo8PAovUmVzb3VyY2VzIDw8Ci9Qcm9jU2V0IFsgL1BERiAvSW1hZ2VDIF0KL1hPYmplY3QgPDwKL2ltYWdlIDEgMCBSCj4+Cj4+Ci9NZWRpYUJveCBbIDAgMCAyMDAuMCAyMDAuMCBdCi9Db250ZW50cyAzIDAgUgovVHlwZSAvUGFnZQovUGFyZW50IDUgMCBSCj4+ZW5kb2JqCjMgMCBvYmo8PAovTGVuZ3RoIDQ3Cj4+c3RyZWFtCnEgMjAwLjAwMDAwMCAwIDAgMjAwLjAwMDAwMCAwIDAgY20gL2ltYWdlIERvIFEKCmVuZHN0cmVhbQplbmRvYmoKNiAwIG9iajw8Ci9DcmVhdGlvbkRhdGUgKEQ6MjAyNjA2MjUwMzM4NDVaKQovTW9kRGF0ZSAoRDoyMDI2MDYyNTAzMzg0NVopCj4+ZW5kb2JqCnhyZWYKMCA3CjAwMDAwMDAwMDAgNjU1MzYgZiAKMDAwMDAwMDE0NCAwMDAwMCBuIAowMDAwMDAxNjE1IDAwMDAwIG4gCjAwMDAwMDE3NzcgMDAwMDAgbiAKMDAwMDAwMDA0MCAwMDAwMCBuIAowMDAwMDAwMDg3IDAwMDAwIG4gCjAwMDAwMDE4NzIgMDAwMDAgbiAKdHJhaWxlcgo8PAovUm9vdCA0IDAgUgovU2l6ZSA3Ci9JbmZvIDYgMCBSCj4+CnN0YXJ0eHJlZgoxOTU0CiUlRU9G"  # noqa: E501


def _png(const: str) -> bytes:
    return base64.b64decode(const)


def _make_pptx(path: Path, blobs: list[bytes]) -> None:
    """构造每页一张图片的 pptx（图片为原生像素尺寸）。"""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for blob in blobs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(BytesIO(blob), Inches(1), Inches(1))
    prs.save(str(path))


def _image(sha: str, blob: bytes | None = None) -> SourceImage:
    return SourceImage(
        page_index=1,
        media_type="image/png",
        blob=blob if blob is not None else _png(_BIG_RED),
        width=160,
        height=160,
        sha256=sha,
    )


# ---- image_extract ----


def test_extract_dedups_and_keeps_page_index(tmp_path):
    path = tmp_path / "deck.pptx"
    _make_pptx(path, [_png(_BIG_RED), _png(_BIG_RED), _png(_BIG_BLUE)])

    images = extract_pptx_images(path, min_side_px=1)

    # 第 1、2 页为同图 → 去重；保留第 1 页与第 3 页
    assert len(images) == 2
    assert images[0].page_index == 1
    assert images[1].page_index == 3
    assert all(img.media_type == "image/png" for img in images)


def test_extract_filters_small_images_by_default(tmp_path):
    path = tmp_path / "small.pptx"
    _make_pptx(path, [_png(_SMALL)])  # 16px < 128 默认阈值

    assert extract_pptx_images(path) == ()


def test_extract_keeps_large_image_by_default(tmp_path):
    path = tmp_path / "big.pptx"
    _make_pptx(path, [_png(_BIG_RED)])

    images = extract_pptx_images(path)
    assert len(images) == 1
    assert images[0].max_side == 160


def test_extract_corrupt_file_returns_empty(tmp_path):
    path = tmp_path / "broken.pptx"
    path.write_bytes(b"not a real pptx")

    assert extract_pptx_images(path) == ()


def test_extract_pdf_images(tmp_path):
    path = tmp_path / "doc.pdf"
    path.write_bytes(base64.b64decode(_PDF_WITH_IMAGE))

    images = extract_pdf_images(path)

    assert len(images) == 1
    assert images[0].page_index == 1
    assert images[0].media_type == "image/jpeg"
    assert images[0].max_side == 200


def test_extract_pdf_corrupt_returns_empty(tmp_path):
    path = tmp_path / "broken.pdf"
    path.write_bytes(b"not a real pdf")

    assert extract_pdf_images(path) == ()


def test_extract_images_dispatches_by_suffix(tmp_path):
    pdf = tmp_path / "doc.pdf"
    pdf.write_bytes(base64.b64decode(_PDF_WITH_IMAGE))
    pptx = tmp_path / "deck.pptx"
    _make_pptx(pptx, [_png(_BIG_RED)])
    docx = tmp_path / "讲义.docx"
    docx.write_bytes(b"irrelevant")

    assert len(extract_images(pdf)) == 1
    assert len(extract_images(pptx)) == 1
    assert extract_images(docx) == ()  # 不支持的类型返回空


# ---- ImageDescriber ----


class _ReplyModel:
    """固定回复的多模态 fake 模型。"""

    def __init__(self, text: str) -> None:
        self.text = text
        self.calls = 0

    async def __call__(self, messages):
        self.calls += 1
        return {"content": [{"type": "text", "text": self.text}]}


class _BoomModel:
    """每次调用都抛错的 fake 模型，用于模拟网络失败。"""

    def __init__(self) -> None:
        self.calls = 0

    async def __call__(self, messages):
        self.calls += 1
        raise RuntimeError("network down")


def test_describe_info_image_returns_text_and_caches(tmp_path):
    model = _ReplyModel("## 产品对比\n- 甲 vs 乙")
    describer = ImageDescriber(model, cache_dir=tmp_path / "cache")

    out = asyncio.run(describer.describe(_image("sha_info")))

    assert out is not None and "产品对比" in out
    assert (tmp_path / "cache" / "sha_info.txt").read_text(encoding="utf-8")


def test_describe_decorative_returns_none_and_caches_empty(tmp_path):
    model = _ReplyModel("装饰图")
    describer = ImageDescriber(model, cache_dir=tmp_path / "cache")

    out = asyncio.run(describer.describe(_image("sha_dec")))

    assert out is None
    # 装饰图缓存为空文件，避免重复识别
    assert (tmp_path / "cache" / "sha_dec.txt").read_text(encoding="utf-8") == ""


def test_describe_too_short_result_is_discarded(tmp_path):
    # 模型异常输出的碎片（如"三个字"）无知识价值，应被丢弃并缓存为空
    model = _ReplyModel("三个字")
    describer = ImageDescriber(model, cache_dir=tmp_path / "cache")

    out = asyncio.run(describer.describe(_image("sha_short")))

    assert out is None
    assert (tmp_path / "cache" / "sha_short.txt").read_text(encoding="utf-8") == ""


def test_describe_cache_hit_skips_model_call(tmp_path):
    cache = tmp_path / "cache"
    cache.mkdir()
    (cache / "sha_hit.txt").write_text("已缓存的转写", encoding="utf-8")
    model = _BoomModel()
    describer = ImageDescriber(model, cache_dir=cache)

    out = asyncio.run(describer.describe(_image("sha_hit")))

    assert out == "已缓存的转写"
    assert model.calls == 0  # 命中缓存不调用模型


def test_describe_failure_returns_none_without_caching(tmp_path):
    model = _BoomModel()
    describer = ImageDescriber(model, cache_dir=tmp_path / "cache")

    out = asyncio.run(describer.describe(_image("sha_fail")))

    assert out is None
    # 失败不写缓存，下次可重试
    assert not (tmp_path / "cache" / "sha_fail.txt").exists()


# ---- enrich ----


class _StubDescriber:
    """鸭子类型的识别器：指定 blob 视为装饰图（返回 None），其余转写。"""

    def __init__(self, decorative_blob: bytes) -> None:
        self._decorative = decorative_blob

    async def describe(self, image: SourceImage):
        if image.blob == self._decorative:
            return None
        return f"第{image.page_index}页转写内容"


def test_enrich_builds_section_with_page_labels(tmp_path):
    path = tmp_path / "deck.pptx"
    _make_pptx(path, [_png(_BIG_RED), _png(_BIG_BLUE)])  # 第1页信息图，第2页装饰图
    stub = _StubDescriber(decorative_blob=_png(_BIG_BLUE))

    out = asyncio.run(enrich_section_with_vision([path], stub))

    assert "## 课件图片信息（视觉识别）" in out
    assert "（第 1 页配图）" in out
    assert "第1页转写内容" in out
    assert "第 2 页" not in out  # 装饰图被略过


def test_enrich_returns_empty_when_no_pptx(tmp_path):
    out = asyncio.run(
        enrich_section_with_vision([tmp_path / "讲义.docx"], _StubDescriber(b""))
    )
    assert out == ""


def test_enrich_returns_empty_when_all_decorative(tmp_path):
    path = tmp_path / "deco.pptx"
    _make_pptx(path, [_png(_BIG_BLUE)])
    stub = _StubDescriber(decorative_blob=_png(_BIG_BLUE))

    out = asyncio.run(enrich_section_with_vision([path], stub))
    assert out == ""
