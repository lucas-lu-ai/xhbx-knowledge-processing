"""解析层与 section_loader 的单元测试。

均使用 ``tmp_path`` 临时生成的合成文件，不依赖真实绩优案例数据。
"""

from __future__ import annotations

from pathlib import Path

from docx import Document
from pptx import Presentation

from insurance_coach_agents.models import FileType
from insurance_coach_agents.parsers import (
    load_section,
    parse_docx,
    parse_pptx,
    parse_txt,
)


def _make_docx(path: Path) -> None:
    doc = Document()
    doc.add_heading("课程主题：测试", level=1)
    doc.add_paragraph("这是正文一。")
    doc.add_heading("1. 一级要点", level=2)
    doc.add_heading("1.1 子要点", level=3)
    doc.add_paragraph("这是正文二。")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "维度"
    table.rows[0].cells[1].text = "信息"
    table.rows[1].cells[0].text = "年龄"
    table.rows[1].cells[1].text = "43岁"
    doc.save(str(path))


def test_parse_docx_preserves_heading_levels(tmp_path):
    # Arrange
    path = tmp_path / "讲义.docx"
    _make_docx(path)

    # Act
    result = parse_docx(path)

    # Assert
    assert result.file_type is FileType.DOCX
    assert "# 课程主题：测试" in result.text
    assert "## 1. 一级要点" in result.text
    assert "### 1.1 子要点" in result.text
    assert result.warnings == ()


def test_parse_docx_converts_table_to_markdown(tmp_path):
    path = tmp_path / "讲义.docx"
    _make_docx(path)

    result = parse_docx(path)

    assert "| 维度 | 信息 |" in result.text
    assert "| --- | --- |" in result.text
    assert "| 年龄 | 43岁 |" in result.text


def test_parse_docx_missing_file_returns_warning(tmp_path):
    result = parse_docx(tmp_path / "不存在.docx")
    assert result.is_empty
    assert result.warnings and "解析失败" in result.warnings[0]


def test_parse_txt_strips_blank_lines_and_whitespace(tmp_path):
    # Arrange
    path = tmp_path / "转写.txt"
    path.write_text("  第一句  \n\n第二句\n   \n第三句", encoding="utf-8")

    # Act
    result = parse_txt(path)

    # Assert
    assert result.text == "第一句\n第二句\n第三句"
    assert result.file_type is FileType.TXT


def test_parse_txt_empty_file_reports_warning(tmp_path):
    path = tmp_path / "空.txt"
    path.write_text("\n\n   \n", encoding="utf-8")

    result = parse_txt(path)

    assert result.is_empty
    assert result.warnings == ("txt 内容为空",)


def test_parse_pptx_extracts_slide_text(tmp_path):
    # Arrange: 用 Title Only 版式放一个标题
    path = tmp_path / "课件.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "第一页标题"
    prs.save(str(path))

    # Act
    result = parse_pptx(path)

    # Assert
    assert result.file_type is FileType.PPTX
    assert "## 第 1 页" in result.text
    assert "第一页标题" in result.text


def test_load_section_dispatches_and_skips_media(tmp_path):
    # Arrange: 构造一个节目录，含 docx + 媒体 + 无关文件
    case_dir = tmp_path / "案例A"
    section_dir = case_dir / "第1节"
    section_dir.mkdir(parents=True)
    _make_docx(section_dir / "讲义.docx")
    (section_dir / "音频.mp3").write_bytes(b"")
    (section_dir / "视频.mp4").write_bytes(b"")
    (section_dir / ".DS_Store").write_bytes(b"")

    # Act
    section = load_section(section_dir)

    # Assert
    assert section.case_name == "案例A"
    assert section.section_name == "第1节"
    assert len(section.files) == 1
    assert section.files[0].file_type is FileType.DOCX
    assert set(section.skipped_media) == {"音频.mp3", "视频.mp4"}


def test_load_section_rejects_non_directory(tmp_path):
    import pytest

    with pytest.raises(NotADirectoryError):
        load_section(tmp_path / "不存在的目录")
