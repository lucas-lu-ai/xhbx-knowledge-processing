"""分组策略与批处理编排的单元测试（fake 模型，不调真实 API）。"""

from __future__ import annotations

import asyncio
import json
from io import BytesIO
from pathlib import Path

from docx import Document
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from insurance_coach_agents.parsers.grouping import (
    group_by_directory,
    group_by_single_file,
    load_group,
)
from insurance_coach_agents.pipeline import run_pipeline


def _make_docx(path: Path, with_content: bool = True) -> None:
    doc = Document()
    if with_content:
        doc.add_heading("课程主题：测试", level=1)
        doc.add_paragraph("正文内容。")
    doc.save(str(path))


def _make_pptx_with_text_and_image(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
    textbox.text_frame.text = "客户需求分析"
    img = Image.new("RGB", (160, 160), "red")
    buffer = BytesIO()
    img.save(buffer, format="PNG")
    buffer.seek(0)
    slide.shapes.add_picture(buffer, Inches(1), Inches(1.5))
    prs.save(str(path))


def _make_cases(tmp_path: Path, n_sections: int = 2) -> Path:
    cases = tmp_path / "cases"
    for i in range(1, n_sections + 1):
        section = cases / "案例A" / f"第{i}节"
        section.mkdir(parents=True)
        _make_docx(section / f"讲义{i}.docx")
        (section / f"音频{i}.mp3").write_bytes(b"")
    return cases


_ASSESS_CONTENT = {
    "worth_storing": True,
    "reason": "包含可复用方法论",
    "topics": ["需求面谈"],
    "serves": {"qa": "high", "recommend": "mid", "exam": "low", "roleplay": "high"},
    "value_score": 0.7,
}
_REVIEW_CONTENT = {
    "passed": True,
    "heading_ok": True,
    "fidelity_ok": True,
    "no_meta_leak": True,
    "issues": [],
    "score": 0.9,
}


class _FakeStructured:
    def __init__(self, content: dict) -> None:
        self.content = content


class _FakeModel:
    async def generate_structured_output(
        self, messages, structured_model, tool_choice=None
    ):
        # 按目标结构返回对应内容：Assessment / ReviewResult。
        if structured_model.__name__ == "ReviewResult":
            return _FakeStructured(_REVIEW_CONTENT)
        return _FakeStructured(_ASSESS_CONTENT)

    async def __call__(self, messages):
        return {"content": [{"type": "text", "text": "# 测试标题\n## 模块\n正文"}]}


def _last_user_content(messages) -> str:
    return str(getattr(messages[-1], "content", ""))


class _CapturingModel(_FakeModel):
    def __init__(self) -> None:
        self.structured_user_materials: list[str] = []
        self.text_user_materials: list[str] = []

    async def generate_structured_output(
        self, messages, structured_model, tool_choice=None
    ):
        self.structured_user_materials.append(_last_user_content(messages))
        return await super().generate_structured_output(
            messages, structured_model, tool_choice
        )

    async def __call__(self, messages):
        self.text_user_materials.append(_last_user_content(messages))
        return await super().__call__(messages)


class _VisionModel:
    async def __call__(self, messages):
        return {
            "content": [
                {"type": "text", "text": "图片显示三步面谈流程和关键话术"}
            ]
        }


class _AutoFixModel:
    def __init__(self, second_review_passes: bool = True) -> None:
        self.review_calls = 0
        self.text_calls = 0
        self.second_review_passes = second_review_passes

    async def generate_structured_output(
        self, messages, structured_model, tool_choice=None
    ):
        if structured_model.__name__ != "ReviewResult":
            return _FakeStructured(_ASSESS_CONTENT)

        self.review_calls += 1
        if self.review_calls == 1:
            return _FakeStructured(
                {
                    "passed": False,
                    "heading_ok": True,
                    "fidelity_ok": False,
                    "no_meta_leak": True,
                    "issues": ["遗漏了关键话术"],
                    "score": 0.5,
                }
            )
        return _FakeStructured(
            {
                "passed": self.second_review_passes,
                "heading_ok": True,
                "fidelity_ok": self.second_review_passes,
                "no_meta_leak": True,
                "issues": [] if self.second_review_passes else ["仍遗漏关键话术"],
                "score": 0.95 if self.second_review_passes else 0.6,
            }
        )

    async def __call__(self, messages):
        self.text_calls += 1
        if self.text_calls == 1:
            return {
                "content": [
                    {"type": "text", "text": "# 测试标题\n## 模块\n初版正文"}
                ]
            }
        return {
            "content": [
                {"type": "text", "text": "# 测试标题\n## 模块\n返修正文\n关键话术"}
            ]
        }


def test_group_by_directory_finds_section_dirs(tmp_path):
    cases = _make_cases(tmp_path, n_sections=2)
    groups = group_by_directory(cases)
    assert len(groups) == 2
    assert all(g.case_name == "案例A" for g in groups)
    # 媒体被记入 skipped_media，docx 进 file_paths
    assert all(len(g.file_paths) == 1 for g in groups)
    assert all(g.skipped_media for g in groups)


def test_group_by_single_file_one_unit_per_file(tmp_path):
    cases = _make_cases(tmp_path, n_sections=2)
    groups = group_by_single_file(cases)
    # 2 个 docx → 2 个单元；mp3 不计入
    assert len(groups) == 2
    assert all(len(g.file_paths) == 1 for g in groups)


def test_load_group_parses_files(tmp_path):
    cases = _make_cases(tmp_path, n_sections=1)
    group = group_by_directory(cases)[0]
    section = load_group(group)
    assert section.primary_text is not None
    assert "课程主题：测试" in section.primary_text


def test_run_pipeline_produces_outputs_and_manifest(tmp_path):
    # Arrange
    cases = _make_cases(tmp_path, n_sections=2)
    out = tmp_path / "out"
    groups = group_by_directory(cases)

    # Act
    results = asyncio.run(run_pipeline(groups, _FakeModel(), concurrency=2, output_dir=out))

    # Assert
    assert len(results) == 2
    assert all(r.status == "ok" for r in results)
    manifest = json.loads((out / "manifest.json").read_text(encoding="utf-8"))
    assert manifest["summary"]["ok"] == 2
    assert manifest["summary"]["worth_storing"] == 2
    assert len(list(out.rglob("*.md"))) == 2
    assert len(list(out.rglob("*.provenance.json"))) == 2


def test_run_pipeline_feeds_page_bound_vision_into_extraction(tmp_path):
    cases = tmp_path / "cases"
    section = cases / "案例A" / "第1节"
    section.mkdir(parents=True)
    _make_pptx_with_text_and_image(section / "课件.pptx")
    out = tmp_path / "out"
    model = _CapturingModel()

    results = asyncio.run(
        run_pipeline(
            group_by_directory(cases),
            model,
            output_dir=out,
            vision=True,
            vision_model=_VisionModel(),
        )
    )

    assert results[0].status == "ok"
    extractor_material = model.text_user_materials[0]
    assert "客户需求分析" in extractor_material
    assert "### 本页配图内容" in extractor_material
    assert "图片显示三步面谈流程和关键话术" in extractor_material
    assert extractor_material.index("客户需求分析") < extractor_material.index(
        "图片显示三步面谈流程和关键话术"
    )


def test_run_pipeline_with_review_records_verdict(tmp_path):
    cases = _make_cases(tmp_path, n_sections=1)
    out = tmp_path / "out"

    results = asyncio.run(
        run_pipeline(
            group_by_directory(cases),
            _FakeModel(),
            output_dir=out,
            vision=False,
            review=True,
        )
    )

    assert results[0].review_passed is True
    assert results[0].review_score == 0.9
    review_reports = list(out.rglob("*.review.md"))
    assert len(review_reports) == 1
    assert "- 是否通过: 是" in review_reports[0].read_text(encoding="utf-8")
    manifest = json.loads((out / "manifest.json").read_text(encoding="utf-8"))
    assert manifest["summary"]["review_failed"] == 0


def test_run_pipeline_auto_fix_revises_failed_review_output(tmp_path):
    cases = _make_cases(tmp_path, n_sections=1)
    out = tmp_path / "out"
    model = _AutoFixModel()

    results = asyncio.run(
        run_pipeline(
            group_by_directory(cases),
            model,
            output_dir=out,
            vision=False,
            review=True,
            auto_fix=True,
        )
    )

    assert model.text_calls == 2
    assert model.review_calls == 2
    assert results[0].review_passed is True
    md = next(path for path in out.rglob("*.md") if not path.name.endswith(".review.md"))
    assert "返修正文" in md.read_text(encoding="utf-8")
    review_report = next(out.rglob("*.review.md")).read_text(encoding="utf-8")
    assert "- 是否通过: 是" in review_report


def test_run_pipeline_review_without_auto_fix_keeps_initial_output(tmp_path):
    cases = _make_cases(tmp_path, n_sections=1)
    out = tmp_path / "out"
    model = _AutoFixModel()

    results = asyncio.run(
        run_pipeline(
            group_by_directory(cases),
            model,
            output_dir=out,
            vision=False,
            review=True,
            auto_fix=False,
        )
    )

    assert model.text_calls == 1
    assert model.review_calls == 1
    assert results[0].review_passed is False
    md = next(path for path in out.rglob("*.md") if not path.name.endswith(".review.md"))
    assert "初版正文" in md.read_text(encoding="utf-8")
    assert "返修正文" not in md.read_text(encoding="utf-8")


def test_run_pipeline_incremental_skips_existing(tmp_path):
    cases = _make_cases(tmp_path, n_sections=2)
    out = tmp_path / "out"
    groups = group_by_directory(cases)

    asyncio.run(run_pipeline(groups, _FakeModel(), output_dir=out))
    # 第二次运行（非 force）应全部跳过
    results = asyncio.run(run_pipeline(groups, _FakeModel(), output_dir=out))
    assert all(r.status == "skipped" for r in results)


def test_run_pipeline_isolates_failures(tmp_path):
    # Arrange: 一个正常节 + 一个空 docx 节（无可解析文本）
    cases = tmp_path / "cases"
    ok_dir = cases / "案例A" / "好节"
    bad_dir = cases / "案例A" / "坏节"
    ok_dir.mkdir(parents=True)
    bad_dir.mkdir(parents=True)
    _make_docx(ok_dir / "好.docx", with_content=True)
    _make_docx(bad_dir / "坏.docx", with_content=False)
    out = tmp_path / "out"

    # Act
    results = asyncio.run(
        run_pipeline(group_by_directory(cases), _FakeModel(), output_dir=out)
    )

    # Assert: 坏节 failed，但好节仍 ok
    by_name = {r.section_name: r for r in results}
    assert by_name["好节"].status == "ok"
    assert by_name["坏节"].status == "failed"
    assert "无可解析" in by_name["坏节"].error
